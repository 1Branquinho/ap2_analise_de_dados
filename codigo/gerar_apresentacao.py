import os
import sys
import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# Configura encoding correto para Windows
sys.stdout.reconfigure(encoding='utf-8')

# Inicializa apresentação widescreen (16:9)
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
blank_layout = prs.slide_layouts[6]  # Layout em branco

# Paleta de Cores Premium (Light Theme - Corporativo)
BG_COLOR = RGBColor(248, 250, 252)        # Slate-50 #F8FAFC
CARD_BG_COLOR = RGBColor(255, 255, 255)   # Branco #FFFFFF
BORDER_COLOR = RGBColor(226, 232, 240)    # Slate-200 #E2E8F0
TEXT_TITLE = RGBColor(15, 23, 42)         # Deep Slate #0F172A
TEXT_LABEL = RGBColor(71, 85, 105)        # Slate-600 #475569
TEXT_BODY = RGBColor(51, 65, 85)          # Slate-700 #334155

# Cores de Destaque (Light Theme)
ACCENT_SKY = RGBColor(2, 132, 199)        # Sky Blue #0284C7
ACCENT_GOLD = RGBColor(217, 119, 6)        # Amber/Gold #D97706
ACCENT_ROSE = RGBColor(225, 29, 72)        # Crimson/Rose #E11D48
ACCENT_EMERALD = RGBColor(5, 150, 105)     # Mint/Green #059669

def apply_background(slide, color=BG_COLOR):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

def create_slide_header(slide, title, category=None):
    apply_background(slide)
    
    # Adiciona a categoria acima do título principal
    if category:
        tx_cat = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.3))
        tf_cat = tx_cat.text_frame
        tf_cat.word_wrap = True
        p_cat = tf_cat.paragraphs[0]
        p_cat.text = category.upper()
        p_cat.font.name = 'Arial'
        p_cat.font.size = Pt(9.5)
        p_cat.font.bold = True
        p_cat.font.color.rgb = ACCENT_SKY
        p_cat.margin_left = p_cat.margin_right = p_cat.margin_top = p_cat.margin_bottom = 0
        
    # Adiciona o título principal
    tx_title = slide.shapes.add_textbox(Inches(0.5), Inches(0.55) if category else Inches(0.4), Inches(12.333), Inches(0.8))
    tf_title = tx_title.text_frame
    tf_title.word_wrap = True
    p_title = tf_title.paragraphs[0]
    p_title.text = title
    p_title.font.name = 'Arial'
    p_title.font.size = Pt(22)
    p_title.font.bold = True
    p_title.font.color.rgb = TEXT_TITLE
    p_title.margin_left = p_title.margin_right = p_title.margin_top = p_title.margin_bottom = 0

def create_card(slide, left, top, width, height, title, subtitle=None, bullets=None, border_color=BORDER_COLOR, bg_color=CARD_BG_COLOR):
    # Adiciona o retângulo de fundo (sem cantos arredondados, visual mais limpo e moderno)
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    shape.line.color.rgb = border_color
    shape.line.width = Pt(1.0)
    
    # Adiciona a caixa de texto
    padding = Inches(0.25)
    txBox = slide.shapes.add_textbox(left + padding, top + padding, width - (2 * padding), height - (2 * padding))
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    
    # Título do Card
    p_title = tf.paragraphs[0]
    p_title.text = title
    p_title.font.name = 'Arial'
    p_title.font.size = Pt(14)
    p_title.font.bold = True
    p_title.font.color.rgb = TEXT_TITLE
    p_title.space_after = Pt(4)
    
    current_p = p_title
    # Subtítulo (opcional)
    if subtitle:
        p_sub = tf.add_paragraph()
        p_sub.text = subtitle
        p_sub.font.name = 'Arial'
        p_sub.font.size = Pt(9.5)
        p_sub.font.italic = True
        p_sub.font.color.rgb = TEXT_LABEL
        p_sub.space_after = Pt(6)
        current_p = p_sub
        
    # Bullet points (opcional)
    if bullets:
        current_p.space_after = Pt(8)
        for idx, b in enumerate(bullets):
            p_bullet = tf.add_paragraph()
            p_bullet.text = b
            p_bullet.font.name = 'Arial'
            p_bullet.font.size = Pt(10.5)
            p_bullet.font.color.rgb = TEXT_BODY
            p_bullet.level = 0
            # Adiciona espaçamento entre os bullets
            p_bullet.space_after = Pt(5)
            
    return shape

def create_table(slide, left, top, width, height, rows, cols, data, col_widths=None):
    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table
    
    if col_widths:
        for idx, w in enumerate(col_widths):
            table.columns[idx].width = w
            
    for r_idx in range(rows):
        for c_idx in range(cols):
            cell = table.cell(r_idx, c_idx)
            cell.text = str(data[r_idx][c_idx])
            
            p = cell.text_frame.paragraphs[0]
            p.font.name = 'Arial'
            p.font.size = Pt(10)
            p.font.color.rgb = RGBColor(255, 255, 255) if r_idx == 0 else TEXT_BODY
            p.font.bold = True if (r_idx == 0 or c_idx == 0) else False
            p.alignment = PP_ALIGN.LEFT if c_idx == 0 else PP_ALIGN.CENTER
            
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            
            # Estiliza o preenchimento de fundo da célula
            cell.fill.solid()
            if r_idx == 0:
                cell.fill.fore_color.rgb = RGBColor(30, 41, 59)  # Slate-800 marinho para cabeçalho
            else:
                if r_idx % 2 == 1:
                    cell.fill.fore_color.rgb = CARD_BG_COLOR  # Branco
                else:
                    cell.fill.fore_color.rgb = RGBColor(241, 245, 249)  # Slate-100 para alternadas
                    
    return table_shape

print("Iniciando a criação dos slides...")

# ════════════════════════════════════════════════════════════════════
# SLIDE 1: Capa
# ════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
apply_background(slide)

# Adiciona faixa de acento no topo
shape_accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.8), Inches(12.333), Inches(0.08))
shape_accent.fill.solid()
shape_accent.fill.fore_color.rgb = ACCENT_SKY
shape_accent.line.fill.background()

# Caixa de texto principal de Capa
tx_capa = slide.shapes.add_textbox(Inches(0.5), Inches(2.2), Inches(12.333), Inches(3.0))
tf = tx_capa.text_frame
tf.word_wrap = True

# Categoria / Tipo de entrega
p0 = tf.paragraphs[0]
p0.text = "APRESENTAÇÃO EXECUTIVA E ANÁLISE CORPORATIVA"
p0.font.name = 'Arial'
p0.font.size = Pt(12)
p0.font.bold = True
p0.font.color.rgb = ACCENT_SKY
p0.space_after = Pt(18)

# Título Principal
p1 = tf.add_paragraph()
p1.text = "Criação de Valor Econômico e Performance de Mercado"
p1.font.name = 'Arial'
p1.font.size = Pt(32)
p1.font.bold = True
p1.font.color.rgb = TEXT_TITLE
p1.space_after = Pt(8)

# Subtítulo
p2 = tf.add_paragraph()
p2.text = "Estudo Comparativo Multicritério de ABEV3, ASAI3 e AMER3 Frente ao IBOV (2021-2025)"
p2.font.name = 'Arial'
p2.font.size = Pt(16)
p2.font.color.rgb = TEXT_LABEL
p2.space_after = Pt(40)

# Informações de Autoria
p3 = tf.add_paragraph()
p3.text = "Elaborado por: Equipe de Análise Contábil e de Mercado"
p3.font.name = 'Arial'
p3.font.size = Pt(11)
p3.font.color.rgb = TEXT_BODY


# ════════════════════════════════════════════════════════════════════
# SLIDE 2: Introdução e Dinâmica Macroeconômica
# ════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
create_slide_header(slide, "Introdução e Dinâmica Macroeconômica (2021-2025)", "Contexto Macro")

bullets_macro = [
    "Juros Altos: Selic elevada impacta captação.",
    "Inflação: Pressão nos custos e no consumo.",
    "Risco: Foco em liquidez e balanços fortes."
]

bullets_pilares = [
    "ABEV3: Consumo defensivo e caixa de R$ 8,4 bi.",
    "ASAI3: Atacarejo alavancado em expansão.",
    "AMER3: Crise de governança e RJ."
]

create_card(slide, Inches(0.5), Inches(1.5), Inches(5.9), Inches(5.3), 
            "O Cenário Econômico Brasileiro", "Pressões, inflação e juros a dois dígitos", bullets_macro)
create_card(slide, Inches(6.933), Inches(1.5), Inches(5.9), Inches(5.3), 
            "Os Três Pilares Corporativos Analisados", "Dinâmicas de capital e setores distintos", bullets_pilares)


# ════════════════════════════════════════════════════════════════════
# SLIDE 3: Perfil Corporativo — Ambev (ABEV3)
# ════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
create_slide_header(slide, "Ambev S.A. (ABEV3) — O Gigante de Consumo Defensivo", "Perfil Corporativo")

bullets_abev_geral = [
    "Liderança: Escala dominante em bebidas.",
    "Distribuição: Ampla capilaridade nacional.",
    "Geração de Caixa: Fluxo forte e recorrente."
]

bullets_abev_patr = [
    "Caixa Líquido: R$ 8,4 bi em caixa (risco zero).",
    "Pricing Power: Poder de repasse de preços."
]

create_card(slide, Inches(0.5), Inches(1.5), Inches(5.9), Inches(5.3), 
            "Visão Geral e Modelo de Negócios", "Escala logística e liderança de mercado", bullets_abev_geral)
create_card(slide, Inches(6.933), Inches(1.5), Inches(5.9), Inches(5.3), 
            "Fortaleza Patrimonial e Operacional", "A proteção do caixa líquido e blindagem contra juros", bullets_abev_patr)


# ════════════════════════════════════════════════════════════════════
# SLIDE 4: Atualidades e Desafios Estratégicos — Ambev + Gráfico
# ════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
create_slide_header(slide, "Ambev (ABEV3) — Premiumização, Digitalização e Gráfico de Preços", "Atualidades e Estratégia")

bullets_abev_at = [
    "Premiumização: Marcas premium elevam margens.",
    "Digitalização: Sucesso com BEES e Zé Delivery.",
    "Riscos Fiscais: Mudanças no JCP e tributos."
]

create_card(slide, Inches(0.5), Inches(1.5), Inches(5.5), Inches(5.3), 
            "Estratégia e Riscos Corporativos", "Premiumização e as novas fronteiras digitais e fiscais", bullets_abev_at)

# Insere o gráfico individual do ticker à direita se existir
img_abev = "graficos/evolucao_abev3.png"
if os.path.exists(img_abev):
    slide.shapes.add_picture(img_abev, Inches(6.3), Inches(1.8), Inches(6.5), Inches(4.5))
else:
    create_card(slide, Inches(6.3), Inches(1.5), Inches(6.5), Inches(5.3), "Gráfico de Preços ABEV3", "Imagem não encontrada no diretório 'graficos/'")


# ════════════════════════════════════════════════════════════════════
# SLIDE 5: Perfil Corporativo — Assaí Atacadista (ASAI3)
# ════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
create_slide_header(slide, "Assaí (ASAI3) — Liderança no Modelo de Atacarejo", "Perfil Corporativo")

bullets_asai_geral = [
    "Foco em Volume: Preços baixos e escala.",
    "Giro Rápido: Alta rotação compensa margem curta.",
    "Independência: Cisão do GPA concluída em 2021."
]

bullets_asai_setor = [
    "Canal Popular: Preferido na inflação alta.",
    "Capital de Giro: Gestão rigorosa de estoques."
]

create_card(slide, Inches(0.5), Inches(1.5), Inches(5.9), Inches(5.3), 
            "O Modelo de Atacarejo", "Giro acelerado, volume e preços agressivos", bullets_asai_geral)
create_card(slide, Inches(6.933), Inches(1.5), Inches(5.9), Inches(5.3), 
            "Dinâmica Operacional do Setor", "O canal preferido do consumidor em tempos de alta inflação", bullets_asai_setor)


# ════════════════════════════════════════════════════════════════════
# SLIDE 6: Atualidades e o Peso do Crescimento — Assaí + Gráfico
# ════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
create_slide_header(slide, "Assaí (ASAI3) — Expansão sob Dívida e Gráfico de Preços", "Atualidades e Estratégia")

bullets_asai_at = [
    "Conversão Extra: Foco em pontos premium.",
    "Expansão: Abertura agressiva de novas lojas.",
    "Alavancagem: Dívida bruta de R$ 26,8 bi.",
    "Juros Altos: Custo financeiro corrói o lucro."
]

create_card(slide, Inches(0.5), Inches(1.5), Inches(5.5), Inches(5.3), 
            "Expansão e Estrutura Financeira", "Conversão de ativos e o custo do carrego da dívida líquida", bullets_asai_at)

# Insere o gráfico individual do ticker à direita se existir
img_asai = "graficos/evolucao_asai3.png"
if os.path.exists(img_asai):
    slide.shapes.add_picture(img_asai, Inches(6.3), Inches(1.8), Inches(6.5), Inches(4.5))
else:
    create_card(slide, Inches(6.3), Inches(1.5), Inches(6.5), Inches(5.3), "Gráfico de Preços ASAI3", "Imagem não encontrada no diretório 'graficos/'")


# ════════════════════════════════════════════════════════════════════
# SLIDE 7: Perfil Corporativo — Americanas (AMER3)
# ════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
create_slide_header(slide, "Americanas (AMER3) — Tradicionalismo e Multicanalidade", "Perfil Corporativo")

bullets_amer_geral = [
    "Presença Física: Lojas em pontos estratégicos.",
    "Multicanalidade: Integração físico e digital.",
    "Perfil: Varejo tradicional de baixo ticket."
]

bullets_amer_vuln = [
    "Pressão Digital: Concorrência agressiva.",
    "Sensibilidade: Dependência de crédito barato."
]

create_card(slide, Inches(0.5), Inches(1.5), Inches(5.9), Inches(5.3), 
            "Visão Histórica e Multicanalidade", "Integração das lojas físicas com o comércio eletrônico", bullets_amer_geral)
create_card(slide, Inches(6.933), Inches(1.5), Inches(5.9), Inches(5.3), 
            "Vulnerabilidades do Modelo", "Concorrência digital agressiva e sensibilidade ao consumo de varejo", bullets_amer_vuln)


# ════════════════════════════════════════════════════════════════════
# SLIDE 8: Atualidades e a Reestruturação — Americanas + Gráfico
# ════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
create_slide_header(slide, "Americanas (AMER3) — O Rombo Contábil e o Gráfico de Preço", "Atualidades e Estratégia")

bullets_amer_at = [
    "Rombo Contábil: R$ 20 bi em risco sacado.",
    "Insolvência: RJ e perda de crédito.",
    "Plano de RJ: Aporte e conversão de dívidas.",
    "Grupamento: Ajuste de 100:1 na B3."
]

create_card(slide, Inches(0.5), Inches(1.5), Inches(5.5), Inches(5.3), 
            "Fraude, Crise de Crédito e RJ", "A reestruturação societária e o grupamento de ações na B3", bullets_amer_at)

# Insere o gráfico individual do ticker à direita se existir
img_amer = "graficos/evolucao_amer3.png"
if os.path.exists(img_amer):
    slide.shapes.add_picture(img_amer, Inches(6.3), Inches(1.8), Inches(6.5), Inches(4.5))
else:
    create_card(slide, Inches(6.3), Inches(1.5), Inches(6.5), Inches(5.3), "Gráfico de Preços AMER3", "Imagem não encontrada no diretório 'graficos/'")


# ════════════════════════════════════════════════════════════════════
# SLIDE 9: Painel Contábil Consolidado (2025)
# ════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
create_slide_header(slide, "Painel Comparativo Contábil e de Desempenho (2025)", "Análise de Balanços")

# Dados da Tabela de Decisão do comparativo_decisao.py (Unidades em R$ milhões, exceto margens)
data_contabil = [
    ["Grupo / Indicador / Métrica", "ABEV3 (Ambev)", "AMER3 (Americanas)", "ASAI3 (Assaí)"],
    ["Receita Líquida", "R$ 82.593 mi", "R$ 12.328 mi", "R$ 71.519 mi"],
    ["Lucro Bruto", "R$ 42.474 mi", "R$ 3.124 mi", "R$ 12.066 mi"],
    ["Margem Bruta (%)", "51.42%", "25.34%", "16.87%"],
    ["EBITDA", "R$ 28.324 mi", "R$ 684 mi", "R$ 3.623 mi"],
    ["Margem EBITDA (%)", "34.29%", "5.55%", "5.07%"],
    ["EBIT (Lucro Operacional)", "R$ 21.917 mi", "R$ 331 mi", "R$ 3.623 mi"],
    ["Margem EBIT (%)", "26.54%", "2.69%", "5.07%"],
    ["Lucro Líquido", "R$ 14.968 mi", "-R$ 271 mi", "R$ 1.023 mi"],
    ["Margem Líquida (%)", "18.12%", "-2.20%", "1.43%"],
    ["Patrimônio Líquido (PL)", "R$ 82.599 mi", "R$ 5.372 mi", "R$ 4.705 mi"],
    ["Dívida Líquida", "-R$ 8.440 mi (Caixa)", "R$ 5.772 mi", "R$ 21.319 mi"]
]

col_widths = [Inches(3.833), Inches(2.833), Inches(2.833), Inches(2.833)]
create_table(slide, Inches(0.5), Inches(1.4), Inches(12.333), Inches(5.4), 12, 4, data_contabil, col_widths)


# ════════════════════════════════════════════════════════════════════
# SLIDE 10: Criação de Valor Econômico (EVA) — 2025
# ════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
create_slide_header(slide, "Economic Value Added (EVA) — Geração vs. Destruição de Valor", "Criação de Valor")

# Explicação conceitual no topo
tx_eva_concept = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(12.333), Inches(0.8))
tf_concept = tx_eva_concept.text_frame
tf_concept.word_wrap = True
tf_concept.margin_left = tf_concept.margin_right = tf_concept.margin_top = tf_concept.margin_bottom = 0
p_concept = tf_concept.paragraphs[0]
p_concept.text = "O EVA mede a riqueza real gerada ao comparar o ROIC contra o WACC. Spread positivo cria valor para o acionista."
p_concept.font.name = 'Arial'
p_concept.font.size = Pt(11)
p_concept.font.color.rgb = TEXT_BODY

# Três cards horizontais para cada empresa
bullets_abev_eva = [
    "ROIC: 19,24% vs. WACC: 14,08%.",
    "Spread: +5,16% (Positivo).",
    "EVA: +R$ 4,1 bilhões.",
    "Geração consistente de valor."
]

bullets_asai_eva = [
    "ROIC: 8,89% vs. WACC: 10,14%.",
    "Spread: -1,25% (Negativo).",
    "EVA: -R$ 333 milhões.",
    "Destruição por custo da dívida."
]

bullets_amer_eva = [
    "ROIC: 2,07% vs. WACC: 12,28%.",
    "Spread: -10,21% (Negativo).",
    "EVA: -R$ 1,07 bilhão.",
    "Destruição severa de capital."
]

# Ambev Card (Green Border for value creation)
create_card(slide, Inches(0.5), Inches(2.2), Inches(3.9), Inches(4.6), 
            "ABEV3 — Ambev", "Criação de Valor Real", bullets_abev_eva, border_color=ACCENT_EMERALD)

# Assaí Card (Orange Border for moderate destruction)
create_card(slide, Inches(4.716), Inches(2.2), Inches(3.9), Inches(4.6), 
            "ASAI3 — Assaí", "Destruição de Valor por Dívida", bullets_asai_eva, border_color=ACCENT_GOLD)

# Americanas Card (Red Border for high destruction)
create_card(slide, Inches(8.933), Inches(2.2), Inches(3.9), Inches(4.6), 
            "AMER3 — Americanas", "Destruição Severa de Riqueza", bullets_amer_eva, border_color=ACCENT_ROSE)


# ════════════════════════════════════════════════════════════════════
# SLIDE 11: Painel de Performance de Mercado (2021-2025)
# ════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
create_slide_header(slide, "Painel de Performance e Risco Acionário Comparado", "Análise de Mercado")

data_mercado = [
    ["Estatística de Mercado (2021-2025)", "ABEV3", "AMER3", "ASAI3", "IBOV (Index)"],
    ["Retorno Acumulado (5 anos)", "29.28%", "-99.94%", "-47.19%", "46.03%"],
    ["Retorno Anualizado (CAGR)", "5.27%", "-76.96%", "-11.99%", "7.87%"],
    ["Volatilidade Anualizada", "23.34%", "129.45%", "40.45%", "17.46%"],
    ["Índice Sharpe (Rf = 10.75%)", "-0.23", "-0.68", "-0.56", "-0.17"],
    ["Beta Sistemático (vs IBOV)", "0.6430", "1.5005", "1.0981", "1.0000"],
    ["Máximo Drawdown (Pior Queda)", "-36.18%", "-99.96%", "-75.55%", "-26.50%"]
]

col_widths_mercado = [Inches(4.333), Inches(2.0), Inches(2.0), Inches(2.0), Inches(2.0)]
create_table(slide, Inches(0.5), Inches(1.5), Inches(12.333), Inches(5.0), 7, 5, data_mercado, col_widths_mercado)


# ════════════════════════════════════════════════════════════════════
# SLIDE 12: Análise Gráfica — Evolução Acumulada vs. IBOV
# ════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
create_slide_header(slide, "Comparativo de Retornos Acumulados (Normalizados)", "Performance de Mercado")

# Gráfico à esquerda
img_comparativo = "graficos/comparativo_retorno.png"
if os.path.exists(img_comparativo):
    slide.shapes.add_picture(img_comparativo, Inches(0.5), Inches(1.8), Inches(6.5), Inches(4.5))
else:
    create_card(slide, Inches(0.5), Inches(1.5), Inches(6.5), Inches(5.3), "Gráfico Comparativo", "Imagem não encontrada no diretório 'graficos/'")

# Comentários à direita
bullets_chart_comment = [
    "IBOV (+46%): Superou as ações do grupo.",
    "ABEV3 (+29%): Ativo defensivo com menor queda.",
    "ASAI3 (-47%): Penalizada por juros e dívida.",
    "AMER3 (-99%): Perda quase total de valor."
]

create_card(slide, Inches(7.333), Inches(1.5), Inches(5.5), Inches(5.3), 
            "Análise do Gráfico de Retornos", "A correlação entre endividamento e desvalorização", bullets_chart_comment)


# ════════════════════════════════════════════════════════════════════
# SLIDE 13: Risco Sistemático (Beta) e Matriz de Correlação
# ════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
create_slide_header(slide, "Beta Sistemático, Correlação e Diversificação", "Performance de Mercado")

# Tabela de Correlação à esquerda
data_corr = [
    ["Correlação", "ABEV3", "AMER3", "ASAI3", "IBOV"],
    ["ABEV3", "1.0000", "0.0878", "0.2511", "0.4809"],
    ["AMER3", "0.0878", "1.0000", "0.1331", "0.2023"],
    ["ASAI3", "0.2511", "0.1331", "1.0000", "0.4738"],
    ["IBOV", "0.4809", "0.2023", "0.4738", "1.0000"]
]

col_widths_corr = [Inches(1.8), Inches(1.1), Inches(1.1), Inches(1.1), Inches(1.1)]
create_table(slide, Inches(0.5), Inches(2.0), Inches(6.2), Inches(3.5), 5, 5, data_corr, col_widths_corr)

# Título da Tabela de Correlação
tx_table_title = slide.shapes.add_textbox(Inches(0.5), Inches(1.4), Inches(6.2), Inches(0.5))
tf_tt = tx_table_title.text_frame
tf_tt.word_wrap = True
tf_tt.margin_left = tf_tt.margin_right = tf_tt.margin_top = tf_tt.margin_bottom = 0
p_tt = tf_tt.paragraphs[0]
p_tt.text = "Matriz de Correlação de Retornos Diários"
p_tt.font.name = 'Arial'
p_tt.font.size = Pt(13)
p_tt.font.bold = True
p_tt.font.color.rgb = TEXT_TITLE

# Comentários à direita
bullets_beta = [
    "ABEV3 (Beta 0,64): Defensiva, reduz risco.",
    "ASAI3 (Beta 1,10): Alta sensibilidade a juros.",
    "AMER3 (Beta 1,50): Oscilações extremas e isoladas."
]

create_card(slide, Inches(7.133), Inches(1.5), Inches(5.7), Inches(5.3), 
            "Beta Sistemático e Alocação", "Risco de mercado e correlação entre ativos", bullets_beta)


# ════════════════════════════════════════════════════════════════════
# SLIDE 14: Simulação de Monte Carlo — Resultados e VaR
# ════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
create_slide_header(slide, "Simulação de Monte Carlo — Projeção Estocástica (2026)", "Projeção de Risco")

# Explicação conceitual no topo
tx_mc_concept = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(12.333), Inches(0.8))
tf_mc = tx_mc_concept.text_frame
tf_mc.word_wrap = True
tf_mc.margin_left = tf_mc.margin_right = tf_mc.margin_top = tf_mc.margin_bottom = 0
p_mc = tf_mc.paragraphs[0]
p_mc.text = "Projeção de 10.000 caminhos de preços para 2026 via Movimento Browniano Geométrico."
p_mc.font.name = 'Arial'
p_mc.font.size = Pt(11)
p_mc.font.color.rgb = TEXT_BODY

# Tabela de Resultados de Monte Carlo
data_mc = [
    ["Projeção Monte Carlo (2026)", "ABEV3", "AMER3", "ASAI3"],
    ["Último Preço Real (Dez/2025)", "R$ 13,86", "R$ 5,13", "R$ 7,20"],
    ["Preço Mediano Projetado", "R$ 14,26", "R$ 1,11", "R$ 6,30"],
    ["Percentil 5% (Pior Cenário)", "R$ 9,67", "R$ 0,14", "R$ 3,23"],
    ["Percentil 95% (Melhor Cenário)", "R$ 20,99", "R$ 9,05", "R$ 12,26"],
    ["VaR 95% (Perda Máxima Projetada)", "-30,27%", "-97,26%", "-55,09%"],
    ["Probabilidade de Perda", "45,1%", "88,6%", "63,0%"],
]

col_widths_mc = [Inches(4.333), Inches(2.666), Inches(2.666), Inches(2.666)]
create_table(slide, Inches(0.5), Inches(2.2), Inches(12.333), Inches(4.5), 7, 4, data_mc, col_widths_mc)


# ════════════════════════════════════════════════════════════════════
# SLIDE 15: Gráficos de Monte Carlo — Funil de Projeção
# ════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
create_slide_header(slide, "Simulação de Monte Carlo — Funis de Projeção de Preços", "Projeção de Risco")

# Insere os 3 gráficos de Monte Carlo lado a lado
mc_images = [
    ("graficos/monte_carlo_abev3.png", "ABEV3"),
    ("graficos/monte_carlo_amer3.png", "AMER3"),
    ("graficos/monte_carlo_asai3.png", "ASAI3"),
]

img_width = Inches(4.0)
img_height = Inches(2.75)
x_positions = [Inches(0.5), Inches(4.666), Inches(8.833)]

for idx, (img_path, ticker) in enumerate(mc_images):
    if os.path.exists(img_path):
        slide.shapes.add_picture(img_path, x_positions[idx], Inches(1.5), img_width, img_height)
    else:
        create_card(slide, x_positions[idx], Inches(1.5), img_width, img_height, 
                    f"Monte Carlo {ticker}", "Imagem não encontrada")

# Card de análise embaixo dos gráficos
bullets_mc_analysis = [
    "ABEV3: Menor risco, VaR de -30,27%.",
    "ASAI3: 63,0% de probabilidade de perda nominal.",
    "AMER3: Risco máximo, 88,6% de chance de perda."
]

create_card(slide, Inches(0.5), Inches(4.5), Inches(12.333), Inches(2.5), 
            "Diagnóstico Estocástico das Projeções", 
            "Interpretação dos funis de probabilidade baseados em 10.000 simulações (MBG)", 
            bullets_mc_analysis)


# ════════════════════════════════════════════════════════════════════
# SLIDE 16: Conclusões e Recomendações Finais
# ════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
create_slide_header(slide, "Conclusões Estratégicas e Recomendações", "Conclusão")

bullets_concl = [
    "Ambev: Caixa líquido e EVA positivo protegem capital.",
    "Assaí: Expansão alavancada sob juros corrói valor.",
    "Americanas: Governança falha gera perdas permanentes."
]

bullets_recom = [
    "Alocação: Focar em geração real de valor (EVA+).",
    "Alavancagem: Evitar dívida alta com Selic a 2 dígitos.",
    "Risco: Monitorar a volatilidade e o VaR projetado."
]

create_card(slide, Inches(0.5), Inches(1.5), Inches(5.9), Inches(5.3), 
            "Conclusões Fundamentais", "Lições obtidas da análise de balanço e mercado", bullets_concl)
create_card(slide, Inches(6.933), Inches(1.5), Inches(5.9), Inches(5.3), 
            "Recomendações Práticas", "Diretrizes de investimento e alocação patrimonial", bullets_recom)

# Salva a apresentação
pptx_filename = "2_Apresentacao_Executiva.pptx"
prs.save(pptx_filename)
print(f"Apresentação salva com sucesso como '{pptx_filename}'.")
